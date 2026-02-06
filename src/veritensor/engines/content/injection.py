# Copyright 2025 Veritensor Security Apache 2.0
# RAG Scanner: Detects Prompt Injections and PII in text, PDF, and Docx files.

import logging
from typing import List, Generator, Set
from pathlib import Path
from veritensor.engines.static.rules import SignatureLoader, is_match
from veritensor.engines.content.pii import PIIScanner  # <--- NEW IMPORT

logger = logging.getLogger(__name__)

# Supported text formats for RAG scanning
TEXT_EXTENSIONS = {
    # Documentation & Markup
    ".txt", ".md", ".markdown", ".rst", ".adoc", ".asciidoc", 
    ".tex", ".org", ".wiki",
    
    # Data & Configs
    ".json", ".xml", ".yaml", ".yml", ".toml", 
    ".ini", ".cfg", ".conf", ".env", ".properties", ".editorconfig",
    
    # Source Code (Scripts)
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".hpp",
    ".rs", ".go", ".rb", ".php", ".pl", ".lua",
    ".sh", ".bash", ".zsh", ".ps1", ".bat", ".sql",
    
    # Infrastructure & DevOps
    ".dockerfile", ".tf", ".tfvars", ".k8s", ".helm", ".tpl",
    ".gitignore", ".gitattributes",
    
    # Logs
    ".log", ".out", ".err"
}

DOC_EXTS = {".pdf", ".docx", ".pptx"}

try:
    import pypdf
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

CHUNK_SIZE = 1024 * 1024 # 1MB chunks
OVERLAP_SIZE = 4096      # 4KB overlap

def scan_document(file_path: Path) -> List[str]:
    """
    Universal entry point for scanning documents (RAG Data).
    Dispatches to specific extractors based on extension.
    """
    ext = file_path.suffix.lower()
    threats = []
    signatures = SignatureLoader.get_prompt_injections()

    try:
        # 1. Get Text Generator (Yields chunks)
        text_generator = None
        
        if ext in TEXT_EXTENSIONS:
            text_generator = _read_text_sliding(file_path)
        elif ext == ".pdf" and PDF_AVAILABLE:
            full_text = _read_pdf(file_path)
            text_generator = _yield_string_chunks(full_text)
        elif ext == ".docx" and DOCX_AVAILABLE:
            full_text = _read_docx(file_path)
            text_generator = _yield_string_chunks(full_text)
        elif ext == ".pptx" and PPTX_AVAILABLE:
            full_text = _extract_text_from_pptx(file_path)
            text_generator = _yield_string_chunks(full_text)
        else:
            return []

        # 2. Scan Chunks
        for chunk_index, chunk in enumerate(text_generator):
            # A. Prompt Injection
            if is_match(chunk, signatures):
                for pattern in signatures:
                    if is_match(chunk, [pattern]):
                        threats.append(f"HIGH: Prompt Injection detected in {file_path.name}: '{pattern}'")
                        return threats # Fail fast

            # B. PII Scan (Presidio) <--- NEW BLOCK
            pii_threats = PIIScanner.scan(chunk)
            if pii_threats:
                threats.extend(pii_threats)
                # We can choose to fail fast on PII too, or collect all.
                # For performance, failing fast on first PII block is often acceptable.
                return threats

    except Exception as e:
        logger.warning(f"Failed to scan document {file_path}: {e}")
        threats.append(f"WARNING: Doc Scan Error: {str(e)}")
        
    return threats

def _read_text_sliding(path: Path) -> Generator[str, None, None]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        buffer = ""
        while True:
            chunk = f.read(CHUNK_SIZE)
            if not chunk:
                break
            data = buffer + chunk
            yield data
            buffer = chunk[-OVERLAP_SIZE:]

def _yield_string_chunks(text: str) -> Generator[str, None, None]:
    if not text: return
    yield text

def _read_pdf(path: Path) -> str:
    text_content = []
    try:
        reader = pypdf.PdfReader(path)
        max_pages = min(len(reader.pages), 50) 
        for i in range(max_pages):
            page_text = reader.pages[i].extract_text()
            if page_text:
                text_content.append(page_text)
        return "\n".join(text_content)
    except Exception as e:
        logger.debug(f"PDF parsing error: {e}")
        return ""

def _read_docx(path: Path) -> str:
    text_content = []
    try:
        doc = docx.Document(path)
        max_paras = min(len(doc.paragraphs), 2000)
        for i in range(max_paras):
            text_content.append(doc.paragraphs[i].text)
        return "\n".join(text_content)
    except Exception as e:
        logger.debug(f"DOCX parsing error: {e}")
        return ""

def _extract_text_from_pptx(path: Path) -> str:
    if not PPTX_AVAILABLE:
        return ""
    text_runs = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
    except Exception as e:
        logger.warning(f"Failed to parse PPTX {path.name}: {e}")
    return "\n".join(text_runs)
